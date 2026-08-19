"""Assemble the generated slide images into a PowerPoint deck.

Run :mod:`make_slides` first, then this, or just run this: it will generate any
missing slide images itself.

Each slide is a full-bleed 16:9 image, because the images already carry their
own titles and captions.  Speaker notes are attached to every slide so the deck
can be presented without a separate script, and so the numbers being quoted are
recorded next to the figure that produced them.

Output: ``results/RFDT_indoor_rf_simulator.pptx``
"""

from __future__ import annotations

import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation                                        # noqa: E402
from pptx.dml.color import RGBColor                                  # noqa: E402
from pptx.util import Inches, Pt                                     # noqa: E402

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
RESULTS = os.path.join(ROOT, "results")
SLIDES = os.path.join(RESULTS, "slides")
OUTPUT = os.path.join(RESULTS, "RFDT_indoor_rf_simulator.pptx")

AUTHOR = "Simran Koul"
TITLE = "A differentiable RF simulator for indoor robotics"
SUBTITLE = ("Transmitter, receiver and material results using RFDT\n"
            "github.com/WIRES-UB/SimranRFDT")

#: Deck order: (image path relative to results/, speaker notes).
#: Set C reuses the experiment figures, which already live in results/.
DECK = [
    ("slides/A1_environment.png",
     "6 x 5 x 2.8 m room, 24 surfaces. Access point on the ceiling at "
     "(1.2, 1.2, 2.7), robot carries the receiver at 0.9 m along an 11.8 m "
     "route. Concrete walls, wood floor, ceiling board, metal cabinet, wooden "
     "table, plasterboard partition. The partition runs from the wall to "
     "y = 3.5 at x = 3 and stands 2.4 m tall, so it shadows the far half of "
     "the room but leaves a gap the robot drives around. Every route is "
     "checked against the furniture; a point inside a solid raises an error."),
    ("slides/A2_ray_trace.png",
     "68 paths found at this position, 14 strongest drawn. Line thickness is "
     "path strength. Top-down shows the bounce points, side view shows how "
     "paths clear obstacles vertically. These are the traced polylines, not a "
     "sketch: the drawn length equals the reported path length exactly."),
    ("slides/A3_route_regimes.png",
     "Ten strongest paths at each of three positions. In line of sight the "
     "direct path carries the link. Near the shadow boundary it weakens. In "
     "deep NLOS behind the partition the channel is carried by reflections "
     "and diffraction. Delay spread and received power quoted per position."),
    ("slides/B1_paths_1.png",
     "One path. The impulse response is a single spike, the amplitude is flat "
     "and the phase is a straight ramp whose slope is exactly the propagation "
     "delay. A single-path channel treats every frequency identically."),
    ("slides/B2_paths_2.png",
     "Two paths. A delayed copy interferes with the direct path and produces "
     "a regular comb of notches. Notch spacing is 1/dtau, annotated on the "
     "slide. This is the same closed form validated to 0.06 percent in the "
     "frequency-response experiment. Selectivity 6.28 dB, delay spread "
     "1.46 ns."),
    ("slides/B3_paths_3.png",
     "Three paths. The comb stops being regular, because three different "
     "delays are now beating against one another. Selectivity 13.05 dB, "
     "delay spread 1.93 ns."),
    ("slides/B4_paths_68.png",
     "All 68 paths. Deep irregular notches, and the phase tears away from the "
     "linear ramp at each one. Selectivity 29.46 dB, delay spread 3.10 ns. "
     "This is what a real indoor channel looks like."),
    ("slides/B5_summary.png",
     "The takeaway. Received power moves by under 3 dB across the whole "
     "sequence, because the direct path dominates it and the direct path "
     "never touches a wall. What multipath changes is the shape of the "
     "channel: frequency selectivity from nothing to 29 dB, delay spread from "
     "zero to 3.1 ns. This is why average signal strength is a poor way to "
     "characterise an indoor link."),
    ("exp5_frequency_response.png",
     "Frequency domain, validated against closed form. Two-ray notch spacing "
     "should be c/dL: 277.49 MHz predicted, 277.67 MHz simulated, 0.06 "
     "percent error. The direct-only response has zero notches and falls "
     "3.522 dB across 4 to 6 GHz, exactly matching the 1/f free-space "
     "prediction. Single-path group delay equals L/c to 1e-15 s. Panel D "
     "measures coherence bandwidth directly and shows the usual 1/(5 sigma) "
     "rule of thumb is fine for metal walls but understates it about "
     "fourfold for concrete and foam."),
    ("exp2_material_sweep.png",
     "Eleven materials from ITU-R P.2040, robot driving the full route, "
     "averaged over the route and a 2 percent frequency band so the result is "
     "not just fading. Delay spread, K-factor and angular spread are strictly "
     "monotone in reflectivity. Metal room behaves like a reverberation "
     "chamber, foam like an anechoic one. At 60 GHz it inverts: reflective "
     "walls raise received power by 16 dB, because the partition blocks the "
     "direct path and bounces are the only energy reaching the far half."),
    ("slides/D1_method_of_images.png",
     "How the second path is found. Mirror the receiver through the wall "
     "plane, join the transmitter to that image with a straight line, and the "
     "bounce point is where the line crosses. Path length equals the "
     "straight-line distance to the image. The detail that matters for this "
     "project: the crossing is computed on the infinite plane, so a bounce "
     "point always exists and moves smoothly when the wall moves. Whether it "
     "lands on the actual wall is applied afterwards as a smooth weight "
     "rather than a yes or no test, which is what keeps the simulator "
     "differentiable with respect to geometry."),
    ("slides/D2_higher_order.png",
     "Third path and beyond. Two bounces: mirror the receiver through both "
     "planes in turn, then sweep forward intersecting each plane. "
     "Diffraction: the bend point on an edge is the one minimising total path "
     "length, which has a closed-form solution, so no iteration is needed and "
     "the gradient stays clean."),
]


def add_title_slide(prs):
    """Opening slide with title, author and the one-line result."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])       # blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x3A)

    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.0),
                                   prs.slide_width - Inches(1.8), Inches(3.4))
    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = TITLE
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    for line in SUBTITLE.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0xA8, 0xC0, 0xD8)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = AUTHOR
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "A differentiable RF ray tracer written from scratch on PyTorch, "
        "implementing the RFDT method (MobiCom 2026), applied to a mobile "
        "robot in a furnished indoor room. Not Sionna: the point of RFDT is "
        "correct gradients with respect to scene geometry, which the paper "
        "shows Sionna does not provide, so the tracer had to be built rather "
        "than wrapped. 29 regression tests validate the physics against "
        "closed-form results: line of sight against Friis to under 0.001 dB, "
        "a two-ray ground reflection to under 0.05 dB, and reciprocity to "
        "1e-13 dB.")
    return slide


def add_image_slide(prs, image_path, notes_text):
    """One full-bleed image slide with speaker notes.

    The images are already 16:9 with their own titles, so they are placed edge
    to edge rather than inside a content placeholder.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(image_path, 0, 0, width=prs.slide_width,
                             height=prs.slide_height)
    slide.notes_slide.notes_text_frame.text = notes_text
    return slide


def add_closing_slide(prs):
    """Final slide: what was validated, and the honest limitations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5),
                                   prs.slide_width - Inches(1.4),
                                   prs.slide_height - Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Validation, and what this does not do"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x2A, 0x3A)

    for heading, items in [
        ("Checked against closed-form physics, 29 tests", [
            "Line of sight against the Friis equation, under 0.001 dB",
            "Two-ray ground reflection against the analytic sum, under 0.05 dB",
            "Reciprocity when transmitter and receiver swap, 1e-13 dB",
            "Two-ray notch spacing against c/dL, 0.06 percent",
            "Autograd against central finite differences, 1e-9 relative",
        ]),
        ("Known limitations", [
            "Scalar polarisation, one complex coefficient per path",
            "First-order diffraction only, none from concave junctions",
            "No measured ground truth: validated against theory, not a "
            "measurement campaign or a full-wave solver",
            "Foam, plastic, paper and human body are approximate literature "
            "values, not ITU data, and are labelled as such in every output",
            "Small scenes: fast for 24 surfaces, not a replacement for a "
            "GPU tracer at 10^5 triangles",
        ]),
    ]:
        p = tf.add_paragraph()
        p.text = ""
        p = tf.add_paragraph()
        p.text = heading
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x2E, 0x86, 0xC1)
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1
            p.font.size = Pt(13.5)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    slide.notes_slide.notes_text_frame.text = (
        "Worth stating the limitations out loud. The validation is against "
        "theory and internal consistency, not against measurements, so the "
        "honest claim is that the physics is implemented correctly, not that "
        "it has been confirmed in a real room.")
    return slide


def ensure_slides():
    """Generate the slide images if they are missing."""
    needed = [os.path.join(RESULTS, rel) for rel, _ in DECK]
    if all(os.path.exists(p) for p in needed):
        return
    print("  slide images missing, generating them first")
    sys.path.insert(0, os.path.join(ROOT, "experiments"))
    __import__("make_slides").main()


def main():
    """Build the deck and write it to results/."""
    ensure_slides()
    prs = Presentation()
    prs.slide_width = Inches(13.333)      # 16:9
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    missing = []
    for rel, notes in DECK:
        path = os.path.join(RESULTS, rel)
        if not os.path.exists(path):
            missing.append(rel)
            continue
        add_image_slide(prs, path, notes)
    add_closing_slide(prs)

    core = prs.core_properties
    core.author = AUTHOR
    core.last_modified_by = AUTHOR
    core.title = TITLE

    prs.save(OUTPUT)
    n_slides = len(prs.slides._sldIdLst)
    print(f"wrote {os.path.relpath(OUTPUT, ROOT)}  "
          f"({n_slides} slides, {os.path.getsize(OUTPUT)/1e6:.1f} MB)")
    if missing:
        print("  missing images, skipped: " + ", ".join(missing))
    print("  every slide carries speaker notes")


if __name__ == "__main__":
    main()